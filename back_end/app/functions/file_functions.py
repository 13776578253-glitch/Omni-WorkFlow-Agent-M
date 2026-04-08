import json, logging
from weasyprint import HTML
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches
from io import BytesIO
from fastapi.responses import StreamingResponse


class PDFFunctions:
    logger = logging.getLogger("file_functions")
    HTML = ""
    HTML_CONTENT = f"""
        <html>
        <head>
        <meta charset="utf-8">
        <style>
        body {{
            font-family: "Noto Sans CJK SC", sans-serif;
        }}
        </style>
        </head>
        <body>
        {HTML}
        </body>
        </html>
    """

    @classmethod
    def html_to_pdf_bytes(cls, html):
        cls.HTML = html
        pdf_io = BytesIO()
        HTML(string=cls.HTML).write_pdf(pdf_io) # 生成PDF写入内存
        pdf_io.seek(0) # 重置指针
        return pdf_io
    
    @classmethod
    def write_pdf(cls, pdf_bytes):
        with open("output.pdf", "wb") as f:
            f.write(pdf_bytes.read())


class PPTFunctions:
    logger = logging.getLogger("file_functions")

    @classmethod
    def generate_ppt_bytes(cls, json_str):
        data = cls.str_to_json(json_str)
        prs = Presentation()

        for slide_data in data["slides"]:
            layout_index = slide_data.get("layout", 1)
            slide_layout = prs.slide_layouts[layout_index]

            slide = prs.slides.add_slide(slide_layout)

            # 标题
            if "title" in slide_data and slide.shapes.title:
                title_shape = slide.shapes.title
                title_shape.text = slide_data["title"]

                cls.apply_style(title_shape, slide_data.get("style", {}), is_title=True)

            # 副标题
            if "sub_title" in slide_data:
                try:
                    subtitle = slide.placeholders[1]
                    subtitle.text = slide_data["sub_title"]
                except:
                    pass

            # 内容
            if "content" in slide_data:
                try:
                    content_shape = cls.get_content_placeholder(slide)

                    if content_shape:
                        tf = content_shape.text_frame
                        tf.clear()  # 清空默认内容

                        content = slide_data["content"]

                        # 如果是list
                        if isinstance(content, list):
                            for i, item in enumerate(content):
                                if i == 0:
                                    p = tf.paragraphs[0]
                                else:
                                    p = tf.add_paragraph()
                                p.text = item
                                p.level = 0  # 一级项目符号

                        # 如果是字符串
                        else:
                            content = content.replace("\\n", "\n")
                            tf.text = content

                        cls.apply_style(content_shape, slide_data.get("style", {}), is_title=False)

                    else:
                        # layout=5 没有内容框则手动创建
                        textbox = slide.shapes.add_textbox(
                            Inches(1), Inches(2), Inches(8), Inches(4)
                        )
                        tf = textbox.text_frame

                        content = slide_data["content"]

                        if isinstance(content, list):
                            for i, item in enumerate(content):
                                if i == 0:
                                    p = tf.paragraphs[0]
                                else:
                                    p = tf.add_paragraph()
                                p.text = item
                        else:
                            tf.text = content.replace("\\n", "\n")

                except Exception as e:
                    cls.logger.error("content error:", e)

        ppt_io = BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)

        return ppt_io

    @classmethod
    def apply_style(cls, shape, style: dict, is_title=False):
        if not shape.has_text_frame: return

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:

                # 字体大小
                if is_title and "font_size" in style:
                    run.font.size = Pt(style["font_size"])
                elif not is_title and "content_font_size" in style:
                    run.font.size = Pt(style["content_font_size"])

                # 字体颜色
                if "font_color" in style:
                    r, g, b = style["font_color"]
                    run.font.color.rgb = RGBColor(r, g, b)

                # 中文字体
                run.font.name = "微软雅黑"

    @classmethod
    def format_content(content: str) -> str:
    # 去掉 markdown 符号
        content = content.replace("**", "")
        return content

    @staticmethod
    def get_content_placeholder(slide):
        for shape in slide.placeholders:
            if shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
                return shape
        return None
    
    @classmethod
    def str_to_json(cls, json_str):
        try:
            data = json.loads(json_str)
        except:
            cls.logger.error("生成json格式有误")
        return data
    
    @classmethod
    def write_ppt(cls, ppt_io):
        with open("output.ppt", "wb") as f:
            f.write(ppt_io.read())

class FileFunction:
    logger = logging.getLogger("file_functions")

    @classmethod
    def pdf_function(cls, data, if_write=None):
        pdf_bytes = PDFFunctions.html_to_pdf_bytes(data)
        if if_write: PDFFunctions.write_pdf(pdf_bytes)
        
        cls.logger.info("已转为pdf文档")
        return pdf_bytes
    
    @classmethod
    def ppt_function(cls, data, if_write=None):
        ppt_bytes = PPTFunctions.generate_ppt_bytes(data)
        if if_write: PPTFunctions.write_ppt(ppt_bytes)
        
        cls.logger.info("已转为ppt文档")
        return ppt_bytes



if __name__ == "__main__":
    from core.logger import setup_logging
    setup_logging()
    html = """<h1>AI生成报告</h1>

    <h2>一、背景</h2>
    <p>这是一个自动生成的PDF文档。</p>

    <h2>二、技术</h2>
    <ul>
    <li>FastAPI</li>
    <li>WeasyPrint</li>
    <li>LLM</li>
    </ul>"""
    json_str = """{
  "slides": [
    {
      "slide_index": 1,
      "layout": 0,
      "title": "React 入门指南",
      "sub_title": "为新手打造的 React 学习路线",
      "style": {
        "font_size": 32,
        "font_color": [0,51,102],
        "align": "center"
      }
    },
    {
      "slide_index": 2,
      "layout": 1,
      "title": "什么是 React?",
      "content": [
        "JavaScript 库",
        "用于构建用户界面的",
        "声明式 UI 框架"
      ],
      "style": {
        "content_font_size": 18
      }
    },
    {
      "slide_index": 3,
      "layout": 1,
      "title": "React 的核心概念",
      "content": [
        "组件 (Components)",
        "JSX (JavaScript XML)",
        "状态 (State)",
        "Props (属性)"
      ],
      "style": {
        "content_font_size": 18
      }
    },
    {
      "slide_index": 4,
      "layout": 5,
      "title": "搭建开发环境",
      "content": [
        "Node.js 和 npm",
        "Create React App"
      ]
    },
    {
      "slide_index": 5,
      "layout": 1,
      "title": "编写你的第一个 React 应用",
      "content": [
        "Hello World",
        "组件的创建和使用",
        "JSX 的语法"
      ],
      "style": {
        "content_font_size": 18
      }
    },
    {
      "slide_index": 6,
      "layout": 5,
      "title": "下一步学习",
      "content": [
        "React Hooks",
        "Redux 或 Context API",
        "React Router"
      ]
    }
  ]
}"""
    
    # FileFunction.pdf_function(html, if_write=True)
    FileFunction.ppt_function(json_str, if_write=True)